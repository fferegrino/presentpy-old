import shlex
from typing import Any, Dict, Iterable, List, Optional, Tuple

import mistletoe
import nbformat
import pygments
import pygments.styles
from mistletoe import block_token
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pygments import lex
from pygments.lexers.python import PythonLexer
from pygments.token import Token

from presentpy.code_cell_config import CodeCellConfig


def get_styles() -> Dict[Any, RGBColor]:
    style = pygments.styles.get_style_by_name("default")
    token_colors = {}
    for token, str_style in style.styles.items():
        if not str_style:
            continue
        _, _, color = str_style.partition("#")
        if not color:
            continue

        pad = 1 if len(color) == 3 else 2
        token_colors[token] = RGBColor(*[int(color[i : i + pad], 16) for i in range(0, len(color), pad)])
    token_colors[Token.Keyword] = RGBColor(0, 128, 20)
    token_colors[Token.Name.Class] = RGBColor(49, 0, 250)
    token_colors[Token.Name.Builtin.Pseudo] = RGBColor(27, 82, 167)
    token_colors[Token.Keyword.Constant] = token_colors[Token.Keyword]
    token_colors[Token.Name.Function.Magic] = token_colors[Token.Name.Class]
    token_colors[Token.Name] = RGBColor(27, 82, 167)
    token_colors[Token.Comment.Single] = RGBColor(76, 135, 135)
    token_colors[Token.Operator] = RGBColor(175, 24, 251)

    return token_colors


token_colors = get_styles()


def add_title_slide(prs: Presentation, title: str, subtitle: Optional[str] = None):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    if subtitle:
        subtitle_shape.text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullet_points: List[str]) -> None:
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    title_shape.text = title

    body_shape = shapes.placeholders[1]

    tf = body_shape.text_frame

    for bullet_point in bullet_points:
        p = tf.add_paragraph()
        p.text = bullet_point
        p.level = 1


def add_code_slide(prs: Presentation, parsed_lines: List[List[Tuple[Any, str]]], config: CodeCellConfig) -> None:
    highlights = config.highlights
    if not highlights:
        highlights = [[0]]
    else:
        highlights = [[0]] + highlights
    for hl in highlights:
        add_code_slide_highlighted(prs, parsed_lines, config.title, highlights=hl)


def add_code_slide_highlighted(
    prs: Presentation, parsed_lines: List[List[Tuple[Any, str]]], title: Optional[str], highlights: Iterable[int]
) -> None:
    highlighted_lines = set(highlights)

    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    if title:
        title_shape = shapes.title
        title_shape.text = title

    body_shape = shapes.placeholders[1]

    text_frame = body_shape.text_frame

    text_frame.clear()

    p = text_frame.paragraphs[0]
    p.bullet = False

    for ln, line in enumerate(parsed_lines, 1):
        for kind, text in line:
            run = p.add_run()
            run.text = text
            font = run.font
            font.bold = ln in highlighted_lines
            font.color.rgb = token_colors.get(kind, RGBColor(0, 0, 0))
            font.name = "Courier"
            font.size = Pt(14)
        run = p.add_run()
        run.font.size = Pt(14)
        run.text = "\x0A"


def get_parsed_lines(source: str) -> List[List[Tuple[Any, str]]]:
    lines = []
    line = []

    for token, value in lex(source, PythonLexer()):
        if token is Token.Text and value == "\n":
            lines.append(line)
            line = []
        else:
            line.append((token, value))

    lines.append(line)

    return lines


def process_notebook(file):
    presentation = Presentation()
    with open(file) as r:
        notebook = nbformat.read(r, as_version=4)

        for cell in notebook["cells"]:
            if cell["cell_type"] == "markdown":
                source = cell["source"]
                document = mistletoe.Document(source)
                header = document.children[0]
                if len(document.children) > 1:
                    if isinstance(document.children[1], block_token.Heading):
                        sub_header = document.children[1]
                        add_title_slide(presentation, header.children[0].content, sub_header.children[0].content)
                    elif isinstance(document.children[1], block_token.List):
                        bullets = [bullet.children[0].children[0].content for bullet in document.children[1].children]
                        add_bullet_slide(
                            presentation,
                            header.children[0].content,
                            bullets,
                        )

            elif cell["cell_type"] == "code":
                source = cell["source"]
                if not source:
                    continue

                source_lines = source.split("\n")
                cell_config = get_config_from_source(source_lines)
                source = "\n".join(source_lines[:-1])
                parsed_lines = get_parsed_lines(source)

                add_code_slide(presentation, parsed_lines, cell_config)
    return presentation


def get_config_from_source(source_lines):
    config = {}
    if source_lines[-1].startswith("#%"):
        config = {
            key: value for key, _, value in [conf.partition("=") for conf in shlex.split(source_lines[-1][2:].strip())]
        }
    dataclass_atrributes = {"title": config.get("title")}

    if highlights := config.get("highlights"):
        lines_to_highlights = highlights.split(",")
        highlight_ints = []
        for l in lines_to_highlights:
            start, _, end = l.partition("-")
            if end:
                highlight_ints.append(list(range(int(start), int(end) + 1)))
            else:
                highlight_ints.append([int(start)])

        dataclass_atrributes["highlights"] = highlight_ints
    cell_config = CodeCellConfig(**dataclass_atrributes)
    return cell_config
